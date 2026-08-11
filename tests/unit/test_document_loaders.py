from docx import Document
from pptx import Presentation

from queryguard.ingestion.docx_loader import extract_docx_units
from queryguard.ingestion.pptx_loader import extract_pptx_units


def test_docx_loader_preserves_heading_locator(tmp_path):
    path = tmp_path / "policy.docx"
    document = Document()
    document.add_heading("Annual Leave", level=1)
    document.add_paragraph("Employees receive 20 days annual leave.")
    document.save(path)

    units, warnings = extract_docx_units(
        path,
        max_uncompressed_bytes=10 * 1024 * 1024,
    )
    assert warnings == []
    assert units[0].locator.startswith("Annual Leave")
    assert "20 days" in units[0].text


def test_pptx_loader_preserves_slide_number(tmp_path):
    path = tmp_path / "review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Risk Summary"
    slide.placeholders[1].text = "Supply chain delay is a major risk."
    presentation.save(path)

    units, warnings = extract_pptx_units(
        path,
        max_uncompressed_bytes=10 * 1024 * 1024,
    )
    assert warnings == []
    assert units[0].locator == "Slide 1"
    assert "Supply chain" in units[0].text
