"""PowerPoint text and table extraction with slide provenance."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from queryguard.documents.models import DocumentUnit
from queryguard.ingestion.common import (
    IngestionError,
    validate_extension,
    validate_office_archive,
)

PPTX_EXTENSIONS = {".pptx"}


def extract_pptx_units(
    path: Path,
    *,
    max_uncompressed_bytes: int,
) -> tuple[list[DocumentUnit], list[str]]:
    validate_extension(path, PPTX_EXTENSIONS)
    validate_office_archive(path, max_uncompressed_bytes=max_uncompressed_bytes)

    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise IngestionError(f"Could not open PPTX {path.name}: {exc}") from exc

    units: list[DocumentUnit] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        table_rows.append(" | ".join(cells))
                if table_rows:
                    blocks.append("\n".join(table_rows))

        text = "\n\n".join(blocks).strip()
        if text:
            units.append(
                DocumentUnit(
                    source_name=path.name,
                    locator=f"Slide {slide_number}",
                    text=text,
                )
            )

    if not units:
        raise IngestionError(f"No readable text was found in {path.name}.")
    return units, []
